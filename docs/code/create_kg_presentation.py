#!/usr/bin/env python3
import os
import io
import sys
from pathlib import Path

# Handle package installation if needed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
except ImportError:
    print("Installing required packages...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    
    # Import after installation
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


class PresentationBuilder:
    """Class to build PowerPoint presentations with consistent styling and layout."""
    
    def __init__(self, widescreen=True):
        """Initialize the presentation with specified dimensions."""
        self.prs = Presentation()
        
        # Set widescreen dimensions if requested
        if widescreen:
            self.prs.slide_width = Inches(13.33)
            self.prs.slide_height = Inches(7.5)
        
        # Define color scheme
        self.colors = {
            'dark_blue': RGBColor(26, 82, 118),      # #1A5276
            'teal': RGBColor(20, 143, 119),          # #148F77
            'orange': RGBColor(230, 126, 34),        # #E67E22
            'white': RGBColor(255, 255, 255),        # #FFFFFF
            'light_gray': RGBColor(240, 240, 240),   # #F0F0F0
            'medium_gray': RGBColor(200, 200, 200),  # #C8C8C8
            'text_dark': RGBColor(50, 50, 50),       # #323232
            'highlight': RGBColor(241, 196, 15)      # #F1C40F
        }
        
        # Store slide references for easy access
        self.slides = []
    
    def add_title_slide(self, title, subtitle):
        """Add a title slide with background styling."""
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        self.slides.append(slide)
        
        # Add background shape with gradient
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, 
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['dark_blue']
        background.line.fill.background()
        
        # Style title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.color.rgb = self.colors['white']
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        
        # Style subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.paragraphs[0].font.color.rgb = self.colors['light_gray']
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        
        return slide
    
    def add_section_slide(self, title, layout_idx=1):
        """Add a standard section slide with consistent styling."""
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)
        self.slides.append(slide)
        
        # Style the title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.color.rgb = self.colors['dark_blue']
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        
        return slide

    def create_text_box(self, slide, left, top, width, height, text="", 
                       font_size=14, font_color=None, bold=False, italic=False,
                       alignment=PP_ALIGN.LEFT, vertical_anchor=MSO_VERTICAL_ANCHOR.TOP):
        """Create a text box with consistent styling."""
        if font_color is None:
            font_color = self.colors['text_dark']
            
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = vertical_anchor
        
        # Add the paragraph and style it
        if text:
            p = text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.bold = bold
            p.font.italic = italic
            p.alignment = alignment
        
        return text_frame
    
    def create_styled_shape(self, slide, shape_type, left, top, width, height, 
                          fill_color=None, line_color=None, text="", font_size=14,
                          font_color=None, bold=False, alignment=PP_ALIGN.CENTER):
        """Create a shape with text and consistent styling."""
        if fill_color is None:
            fill_color = self.colors['light_gray']
        if line_color is None:
            line_color = self.colors['dark_blue']
        if font_color is None:
            font_color = self.colors['text_dark']
            
        # Create the shape
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        
        # Style the outline
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
        
        # Add text if provided
        if text:
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.text = text
            
            # Style the text
            p = text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.bold = bold
            p.alignment = alignment
        
        return shape
    
    def add_bullet_list(self, text_frame, items, font_size=14, font_color=None, level=0):
        """Add a bullet list to an existing text frame."""
        if font_color is None:
            font_color = self.colors['text_dark']
            
        for i, item in enumerate(items):
            # For first item, use the existing paragraph
            if i == 0 and not text_frame.text:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.level = level
        
        return text_frame
    
    def create_comparison_columns(self, slide, titles, contents, fill_colors=None):
        """Create a two-column comparison with titles and bullet points."""
        if fill_colors is None:
            fill_colors = [self.colors['light_gray'], self.colors['teal']]
        
        # Calculate positions
        margin = Inches(1)
        column_spacing = Inches(0.5)
        column_width = (self.prs.slide_width - (2 * margin) - column_spacing) / 2
        top = Inches(2)
        height = Inches(4)
        
        columns = []
        
        for i in range(2):
            # Calculate column position
            left = margin + (i * (column_width + column_spacing))
            
            # Create column background
            box = self.create_styled_shape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, column_width, height,
                fill_color=fill_colors[i],
                line_color=self.colors['dark_blue']
            )
            
            # Add column title
            title_top = top + Inches(0.5)
            title_frame = self.create_text_box(
                slide, left + Inches(0.5), title_top, 
                column_width - Inches(1), Inches(0.5),
                text=titles[i], font_size=20, bold=True,
                font_color=self.colors['dark_blue'] if i == 0 else self.colors['white']
            )
            
            # Add content as bullet points
            content_top = title_top + Inches(0.7)
            content_frame = self.create_text_box(
                slide, left + Inches(0.5), content_top,
                column_width - Inches(1), height - Inches(1.2)
            )
            
            self.add_bullet_list(
                content_frame, contents[i], font_size=16,
                font_color=self.colors['text_dark'] if i == 0 else self.colors['white']
            )
            
            columns.append((box, title_frame, content_frame))
        
        return columns
    
    def create_process_flow(self, slide, steps, descriptions=None, horizontal=True):
        """Create a process flow diagram with steps and optional descriptions."""
        if descriptions is None:
            descriptions = [""] * len(steps)
        
        # Calculate dimensions and positions
        if horizontal:
            # Horizontal flow
            box_width = Inches(2)
            box_height = Inches(1.2)
            arrow_length = Inches(1)
            
            total_width = (box_width * len(steps)) + (arrow_length * (len(steps) - 1))
            start_left = (self.prs.slide_width - total_width) / 2
            top = Inches(2.5)
            
            for i, (step, desc) in enumerate(zip(steps, descriptions)):
                # Calculate positions
                left = start_left + (i * (box_width + arrow_length))
                
                # Add step box
                box = self.create_styled_shape(
                    slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, top, box_width, box_height,
                    fill_color=self.colors['teal'],
                    line_color=self.colors['dark_blue'],
                    text=step, font_size=14, bold=True,
                    font_color=self.colors['white']
                )
                
                # Add description if provided
                if desc:
                    desc_frame = self.create_text_box(
                        slide, left, top + Inches(0.6), 
                        box_width, Inches(0.6),
                        text=desc, font_size=10,
                        font_color=self.colors['white'],
                        alignment=PP_ALIGN.CENTER
                    )
                
                # Add arrow (except for last step)
                if i < len(steps) - 1:
                    arrow_left = left + box_width
                    arrow_top = top + (box_height / 2) - Inches(0.1)
                    
                    arrow = self.create_styled_shape(
                        slide, MSO_SHAPE.RIGHT_ARROW,
                        arrow_left, arrow_top, 
                        arrow_length, Inches(0.2),
                        fill_color=self.colors['orange'],
                        line_color=None
                    )
        else:
            # Vertical flow
            box_width = Inches(10)
            box_height = Inches(0.7)
            arrow_height = Inches(0.4)
            
            total_height = (box_height * len(steps)) + (arrow_height * (len(steps) - 1))
            left = (self.prs.slide_width - box_width) / 2
            start_top = Inches(1.8)
            
            for i, (step, desc) in enumerate(zip(steps, descriptions)):
                # Calculate positions
                top = start_top + (i * (box_height + arrow_height))
                
                # Alternate colors
                fill_color = self.colors['light_gray'] if i % 2 == 0 else self.colors['teal']
                font_color = self.colors['text_dark'] if i % 2 == 0 else self.colors['white']
                
                # Add step box
                box = self.create_styled_shape(
                    slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, top, box_width, box_height,
                    fill_color=fill_color,
                    line_color=self.colors['dark_blue']
                )
                
                # Add step number in circle
                circle_size = Inches(0.5)
                circle_left = left + Inches(0.1)
                circle_top = top + (box_height - circle_size) / 2
                
                circle = self.create_styled_shape(
                    slide, MSO_SHAPE.OVAL,
                    circle_left, circle_top, 
                    circle_size, circle_size,
                    fill_color=self.colors['orange'],
                    line_color=None,
                    text=str(i + 1), font_size=14, 
                    bold=True, font_color=self.colors['white']
                )
                
                # Add step text
                text_frame = self.create_text_box(
                    slide, left + Inches(0.7), top + Inches(0.1),
                    box_width - Inches(0.8), box_height - Inches(0.2),
                    text=step, font_size=16,
                    font_color=font_color
                )
                
                # Add arrow (except for last step)
                if i < len(steps) - 1:
                    arrow_left = left + (box_width / 2) - Inches(0.1)
                    arrow_top = top + box_height
                    
                    arrow = self.create_styled_shape(
                        slide, MSO_SHAPE.DOWN_ARROW,
                        arrow_left, arrow_top, 
                        Inches(0.2), arrow_height,
                        fill_color=self.colors['orange'],
                        line_color=None
                    )
    
    def create_layered_architecture(self, slide, layers, descriptions=None):
        """Create a layered architecture diagram."""
        if descriptions is None:
            descriptions = [""] * len(layers)
            
        # Calculate dimensions and positions
        layer_height = Inches(0.8)
        width = Inches(10)
        left = (self.prs.slide_width - width) / 2
        start_top = Inches(1.8)
        
        for i, (layer_name, description, color) in enumerate(layers):
            top = start_top + (i * layer_height)
            
            # Create layer box
            layer_box = self.create_styled_shape(
                slide, MSO_SHAPE.RECTANGLE,
                left, top, width, layer_height,
                fill_color=color,
                line_color=None
            )
            
            # Add layer name
            name_frame = self.create_text_box(
                slide, left + Inches(0.2), top + Inches(0.1),
                Inches(3), Inches(0.6),
                text=layer_name, font_size=16, bold=True,
                font_color=self.colors['white'] if color in [self.colors['dark_blue'], self.colors['teal']] else self.colors['text_dark']
            )
            
            # Add layer description
            if description:
                desc_frame = self.create_text_box(
                    slide, left + Inches(3.5), top + Inches(0.1),
                    Inches(6.5), Inches(0.6),
                    text=description, font_size=14,
                    font_color=self.colors['white'] if color in [self.colors['dark_blue'], self.colors['teal']] else self.colors['text_dark']
                )
    
    def create_chart(self, slide, chart_type, categories, series_data, 
                    left=None, top=None, width=None, height=None, 
                    has_legend=True, title=None):
        """Create a chart with proper styling."""
        # Set default dimensions if not specified
        if left is None:
            left = Inches(2)
        if top is None:
            top = Inches(2)
        if width is None:
            width = Inches(9)
        if height is None:
            height = Inches(4.5)
        
        # Create the chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        # Add each data series
        for series_name, values in series_data:
            chart_data.add_series(series_name, values)
        
        # Add the chart to the slide
        chart = slide.shapes.add_chart(
            chart_type, left, top, width, height, chart_data
        ).chart
        
        # Add title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        
        # Configure legend
        chart.has_legend = has_legend
        if has_legend:
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
        
        # Style the chart
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.font.size = Pt(12)
        chart.plots[0].data_labels.font.color.rgb = self.colors['text_dark']
        
        if chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
            chart.plots[0].data_labels.number_format = '0"%"'
        
        # Style the axes
        chart.category_axis.has_major_gridlines = False
        chart.value_axis.has_major_gridlines = True
        
        return chart
    
    def save(self, filename):
        """Save the presentation to a file."""
        # Ensure directory exists
        output_dir = os.path.dirname(filename)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
        self.prs.save(filename)
        print(f"Presentation saved successfully to: {filename}")
        return filename


def create_kg_presentation():
    """Create a PowerPoint presentation on Knowledge Graphs for Root Cause Analysis."""
    
    # Create presentation builder
    builder = PresentationBuilder(widescreen=True)
    
    # ----- Slide 1: Title Slide -----
    title_slide = builder.add_title_slide(
        title="Knowledge Graphs: Transforming Root Cause Analysis",
        subtitle="A Systematic Approach to Discovering Hidden Connections"
    )
    
    # ----- Slide 2: The Challenge of Root Cause Analysis -----
    challenge_slide = builder.add_section_slide("The Challenge of Root Cause Analysis")
    
    # Create comparison columns
    builder.create_comparison_columns(
        challenge_slide,
        titles=["Traditional Troubleshooting", "Knowledge Graph Approach"],
        contents=[
            [
                "Siloed information across systems",
                "Hidden relationships between components",
                "Time-consuming manual investigation",
                "Dependency on expert knowledge",
                "Difficulty scaling across systems"
            ],
            [
                "Unified view of interconnected components",
                "Explicit representation of relationships",
                "Systematic traversal of causal paths",
                "Augmented expertise with graph analytics",
                "Reproducible, data-driven insights"
            ]
        ],
        fill_colors=[builder.colors['light_gray'], builder.colors['teal']]
    )
    
    # ----- Slide 3: Knowledge Graph Fundamentals -----
    fundamentals_slide = builder.add_section_slide("Knowledge Graph Fundamentals")
    
    # Create a visual placeholder for the graph
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(4)
    
    # Placeholder shape for image
    kg_viz = builder.create_styled_shape(
        fundamentals_slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
        fill_color=builder.colors['light_gray'],
        line_color=builder.colors['dark_blue']
    )
    
    # Add visualization label
    viz_label = builder.create_text_box(
        fundamentals_slide, left + Inches(2), top + Inches(1.5),
        Inches(2), Inches(1),
        text="Knowledge Graph Visualization", 
        font_size=16, bold=True,
        alignment=PP_ALIGN.CENTER
    )
    
    # Add key components
    components_frame = builder.create_text_box(
        fundamentals_slide, left + width + Inches(0.5), top,
        Inches(4.5), height
    )
    
    components = [
        ("Nodes (Entities)", "Represent components, states, or events"),
        ("Edges (Relationships)", "Connect entities with semantic meaning"),
        ("Properties", "Attributes providing context and details"),
        ("Ontology", "Structured vocabulary defining entity types")
    ]
    
    for component, description in components:
        p = components_frame.add_paragraph()
        p.text = component
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = builder.colors['dark_blue']
        
        p = components_frame.add_paragraph()
        p.text = description
        p.font.size = Pt(16)
        p.space_after = Pt(12)
    
    # ----- Slide 4: Knowledge Graph Architecture -----
    architecture_slide = builder.add_section_slide("Knowledge Graph Architecture for RCA")
    
    # Create layered architecture
    layers = [
        ("Data Acquisition Layer", "Maintenance records, sensor data, equipment specs", builder.colors['light_gray']),
        ("Knowledge Extraction", "Entity and relationship identification", builder.colors['orange']),
        ("Knowledge Representation", "Triple storage with temporal aspects", builder.colors['teal']),
        ("Inference Engine", "Causal pattern detection algorithms", builder.colors['orange']),
        ("Query Interface", "Traversal and visualization capabilities", builder.colors['dark_blue'])
    ]
    
    builder.create_layered_architecture(architecture_slide, layers)
    
    # ----- Slide 5: Graph Data Model -----
    data_model_slide = builder.add_section_slide("Graph Data Model for Failure Analysis")
    
    # Create content
    left_margin = Inches(1.5)
    top_margin = Inches(1.8)
    
    # Entity Types Column
    entity_width = Inches(4.5)
    entity_height = Inches(4.5)
    
    entity_box = builder.create_styled_shape(
        data_model_slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left_margin, top_margin, entity_width, entity_height,
        fill_color=builder.colors['light_gray'],
        line_color=builder.colors['dark_blue']
    )
    
    # Entity Types Header
    entity_header = builder.create_text_box(
        data_model_slide, left_margin + Inches(0.5), top_margin + Inches(0.2),
        entity_width - Inches(1), Inches(0.5),
        text="Entity Types", font_size=20, bold=True,
        font_color=builder.colors['dark_blue']
    )
    
    # Entity Types Content
    entity_types = [
        ("Physical Components", "Equipment, parts, systems"),
        ("Operational States", "Normal, abnormal, failure modes"),
        ("Events", "Failures, maintenance activities"),
        ("Environmental Context", "Operating conditions, settings")
    ]
    
    content_top = top_margin + Inches(0.8)
    for i, (entity_type, entity_desc) in enumerate(entity_types):
        top = content_top + (i * Inches(0.8))
        
        type_frame = builder.create_text_box(
            data_model_slide, left_margin + Inches(0.5), top,
            entity_width - Inches(1), Inches(0.8)
        )
        
        p = type_frame.add_paragraph()
        p.text = entity_type
        p.font.bold = True
        p.font.size = Pt(16)
        
        p = type_frame.add_paragraph()
        p.text = entity_desc
        p.font.size = Pt(14)
    
    # Relationship Types Column
    rel_left = left_margin + entity_width + Inches(0.5)
    rel_width = entity_width
    rel_height = entity_height
    
    rel_box = builder.create_styled_shape(
        data_model_slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        rel_left, top_margin, rel_width, rel_height,
        fill_color=builder.colors['teal'],
        line_color=builder.colors['dark_blue']
    )
    
    # Relationship Types Header
    rel_header = builder.create_text_box(
        data_model_slide, rel_left + Inches(0.5), top_margin + Inches(0.2),
        rel_width - Inches(1), Inches(0.5),
        text="Relationship Types", font_size=20, bold=True,
        font_color=builder.colors['white']
    )
    
    # Relationship Type Content
    relationship_types = [
        ("Causal", "causes, contributes_to"),
        ("Structural", "contains, connects_to"),
        ("Temporal", "precedes, during"),
        ("Functional", "monitors, controls")
    ]
    
    content_top = top_margin + Inches(0.8)
    for i, (rel_type, rel_desc) in enumerate(relationship_types):
        top = content_top + (i * Inches(0.8))
        
        type_frame = builder.create_text_box(
            data_model_slide, rel_left + Inches(0.5), top,
            rel_width - Inches(1), Inches(0.8)
        )
        
        p = type_frame.add_paragraph()
        p.text = rel_type
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = builder.colors['white']
        
        p = type_frame.add_paragraph()
        p.text = rel_desc
        p.font.size = Pt(14)
        p.font.color.rgb = builder.colors['white']
    
    # ----- Slide 6: Building the Knowledge Graph -----
    building_slide = builder.add_section_slide("Building the Knowledge Graph")
    
    # Create process flow
    stages = [
        "Data Integration",
        "Entity Extraction",
        "Relationship Detection",
        "Causal Inference",
        "Knowledge Enhancement"
    ]
    
    descriptions = [
        "Connect maintenance systems, IoT sensors, manuals",
        "Identify components, states, and events",
        "Establish connections between entities",
        "Apply rules and statistical analysis",
        "Enrich with domain expertise"
    ]
    
    builder.create_process_flow(building_slide, stages, descriptions)
    
    # ----- Slide 7: Case Study -----
    case_study_slide = builder.add_section_slide("Case Study - Manufacturing Equipment")
    
    # Left side: placeholder for knowledge graph visualization
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(6)
    height = Inches(4.5)
    
    # Placeholder for knowledge graph visualization
    case_viz = builder.create_styled_shape(
        case_study_slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
        fill_color=builder.colors['light_gray'],
        line_color=builder.colors['dark_blue']
    )
    
    # Add visualization label
    case_label = builder.create_text_box(
        case_study_slide, left + Inches(1.5), top + Inches(2),
        Inches(3), Inches(0.5),
        text="Pump Failure Knowledge Graph", 
        font_size=18, bold=True,
        alignment=PP_ALIGN.CENTER
    )
    
    # Right side: Case study details
    right_left = left + width + Inches(0.5)
    right_width = Inches(4.5)
    
    # Scenario box
    scenario_top = top
    scenario_height = Inches(1)
    
    scenario_box = builder.create_styled_shape(
        case_study_slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, scenario_top, right_width, scenario_height,
        fill_color=builder.colors['dark_blue'],
        line_color=None,
        text="Scenario: Recurring pump bearing failures",
        font_size=16, bold=True, font_color=builder.colors['white'],
        alignment=PP_ALIGN.CENTER
    )
    
    # Approach box
    approach_top = scenario_top + scenario_height + Inches(0.2)
    approach_height = Inches(2.1)
    
    approach_box = builder.create_styled_shape(
        case_study_slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, approach_top, right_width, approach_height,
        fill_color=builder.colors['teal'],
        line_color=None
    )
    
    # Approach text
    approach_frame = builder.create_text_box(
        case_study_slide, right_left + Inches(0.2), approach_top + Inches(0.2),
        right_width - Inches(0.4), approach_height - Inches(0.4)
    )
    
    p = approach_frame.add_paragraph()
    p.text = "Knowledge Graph Approach:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = builder.colors['white']
    
    builder.add_bullet_list(
        approach_frame,
        [
            "Integrated 3 years of maintenance records",
            "Connected 1,200+ sensor readings",
            "Mapped component relationships",
            "Applied causal pattern detection"
        ],
        font_size=12,
        font_color=builder.colors['white']
    )
    
    # Results box
    results_top = approach_top + approach_height + Inches(0.2)
    results_height = Inches(1.8)
    
    results_box = builder.create_styled_shape(
        case_study_slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        right_left, results_top, right_width, results_height,
        fill_color=builder.colors['orange'],
        line_color=None
    )
    
    # Results text
    results_frame = builder.create_text_box(
        case_study_slide, right_left + Inches(0.2), results_top + Inches(0.2),
        right_width - Inches(0.4), results_height - Inches(0.4)
    )
    
    p = results_frame.add_paragraph()
    p.text = "Results:"
    p.font.bold = True
    p.font.size = Pt(14)
    
    builder.add_bullet_list(
        results_frame,
        [
            "Root cause identified: Shaft misalignment",
            "85% reduction in failures",
            "$145,000 annual savings",
            "Implementation time: 6 weeks"
        ],
        font_size=12
    )
    
    # ----- Slide 8: Implementation Roadmap -----
    roadmap_slide = builder.add_section_slide("Implementation Roadmap")
    
    # Add implementation steps
    phases = [
        ("Phase 1", "4 weeks", "Data source integration and schema design"),
        ("Phase 2", "6 weeks", "Knowledge extraction and graph construction"),
        ("Phase 3", "3 weeks", "Causal inference rule development"),
        ("Phase 4", "5 weeks", "Query interface and visualization dashboard"),
        ("Phase 5", "ongoing", "Continuous refinement and expansion")
    ]
    
    # Create a visualization showing the phases
    timeline_top = Inches(2)
    timeline_left = Inches(1.5)
    timeline_width = Inches(10)
    timeline_height = Inches(0.1)
    
    # Timeline line
    timeline = builder.create_styled_shape(
        roadmap_slide, MSO_SHAPE.RECTANGLE,
        timeline_left, timeline_top, timeline_width, timeline_height,
        fill_color=builder.colors['dark_blue'],
        line_color=None
    )
    
    # Add phases
    phase_width = Inches(1.8)
    phase_height = Inches(2.5)
    
    for i, (phase, duration, description) in enumerate(phases):
        left = timeline_left + (i * (timeline_width / len(phases)))
        
        # Add circle marker on timeline
        marker_size = Inches(0.3)
        marker_left = left + ((timeline_width / len(phases)) / 2) - (marker_size / 2)
        marker_top = timeline_top - (marker_size / 2) + (timeline_height / 2)
        
        marker = builder.create_styled_shape(
            roadmap_slide, MSO_SHAPE.OVAL,
            marker_left, marker_top, marker_size, marker_size,
            fill_color=builder.colors['orange'],
            line_color=None
        )
        
        # Add phase info - alternate above and below timeline
        box_left = left + Inches(0.2)
        if i % 2 == 0:
            box_top = timeline_top + Inches(0.5)
        else:
            box_top = timeline_top - phase_height - Inches(0.4)
        
        # Phase box
        phase_box = builder.create_styled_shape(
            roadmap_slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            box_left, box_top, phase_width, phase_height,
            fill_color=builder.colors['light_gray'] if i % 2 == 0 else builder.colors['teal'],
            line_color=None
        )
        
        # Phase text
        text_color = builder.colors['dark_blue'] if i % 2 == 0 else builder.colors['white']
        
        # Phase title
        title_frame = builder.create_text_box(
            roadmap_slide, box_left + Inches(0.1), box_top + Inches(0.1),
            phase_width - Inches(0.2), Inches(0.4),
            text=phase, font_size=16, bold=True,
            font_color=text_color
        )
        
        # Phase duration
        duration_frame = builder.create_text_box(
            roadmap_slide, box_left + Inches(0.1), box_top + Inches(0.5),
            phase_width - Inches(0.2), Inches(0.4),
            text=duration, font_size=14, italic=True,
            font_color=text_color
        )
        
        # Phase description
        desc_frame = builder.create_text_box(
            roadmap_slide, box_left + Inches(0.1), box_top + Inches(1),
            phase_width - Inches(0.2), Inches(1.4),
            text=description, font_size=12,
            font_color=text_color
        )
    
    # ----- Slide 9: Key Performance Metrics -----
    metrics_slide = builder.add_section_slide("Key Performance Metrics")
    
    # Create a performance chart with normalization
    # Data represents: before (base 100) and after (relative change)
    chart = builder.create_chart(
        metrics_slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        ['MTTR', 'First-time Fix', 'Recurring Issues', 'Predictive Ratio'],
        [
            ('Before', (100, 100, 100, 30)),
            ('After', (35, 142, 22, 70))
        ],
        title="Performance Improvement with Knowledge Graph Analytics"
    )
    
    # Add explanatory notes
    notes_frame = builder.create_text_box(
        metrics_slide, Inches(2), Inches(6.5),
        Inches(9), Inches(0.7)
    )
    
    p = notes_frame.add_paragraph()
    p.text = "Note: Values are normalized with 'Before' state at baseline 100% (except Predictive Ratio). Lower values for MTTR and Recurring Issues indicate improvement."
    p.font.size = Pt(12)
    p.font.italic = True
    
    # ----- Slide 10: Getting Started -----
    getting_started_slide = builder.add_section_slide("Getting Started")
    
    # Create implementation steps
    steps = [
        "Define failure modes of interest",
        "Inventory available data sources",
        "Create minimum viable knowledge model",
        "Implement core causal patterns",
        "Validate with historical cases",
        "Expand incrementally"
    ]
    
    # Create a vertical process flow with the steps
    builder.create_process_flow(getting_started_slide, steps, horizontal=False)
    
    # ----- Slide 11: Q&A -----
    qa_slide = builder.add_section_slide("Questions & Discussion")
    
    # Add a large question mark icon
    q_size = Inches(4)
    q_left = (builder.prs.slide_width - q_size) / 2
    q_top = (builder.prs.slide_height - q_size) / 2
    
    q_mark = builder.create_styled_shape(
        qa_slide, MSO_SHAPE.OVAL,
        q_left, q_top, q_size, q_size,
        fill_color=builder.colors['teal'],
        line_color=None
    )
    
    # Add question mark text
    q_text = builder.create_text_box(
        qa_slide, q_left, q_top,
        q_size, q_size,
        text="?", font_size=150, bold=True,
        font_color=builder.colors['white'],
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    )
    
    # Add contact info
    contact_frame = builder.create_text_box(
        qa_slide, Inches(3), Inches(6),
        Inches(7), Inches(0.5),
        text="Contact: your.email@example.com | LinkedIn: yourname",
        font_size=14, alignment=PP_ALIGN.CENTER
    )
    
    # Save the presentation
    output_path = "Knowledge_Graph_for_Root_Cause_Analysis.pptx"
    builder.save(output_path)
    
    return output_path

if __name__ == "__main__":
    create_kg_presentation()